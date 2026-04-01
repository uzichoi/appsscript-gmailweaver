# src/app.py

import datetime
import os
import re
import subprocess
import time
import sys
import json
import threading
import uuid
import openai  
import base64
import requests
import shutil
import zlib

from dotenv import load_dotenv
from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS
import fitz  # PyMuPDF
from docx import Document
import olefile
import csv
from pptx import Presentation
from openpyxl import load_workbook
from flask import send_from_directory

# Job 이용 공통함수 import
from util.jobs.job_store import *
from util.jobs.job_run import start_graph_pipeline_background, start_graph_update_pipeline_background
from config.settings import *
from util.user_path import UserPaths

# 환경변수 로드
load_dotenv("src/parquet/.env") # src/parquet/.env를 사용하는 이유: GraphRAG 설정(settings.yaml)과 API 키가 같은 디렉터리에 위치하기 때문

# Flask 앱 초기화
app = Flask(__name__)   # Flask 앱 객체 생성. 해당 파일이 서버의 메인 애플리케이션이라는 의미
CORS(app)   # Cross-Origin Resource Sharing 허용 (다른 환경에서 이 서버의 API를 호출할 수 있도록)

# Apps Script Web App URL (캘린더, 라벨 등 모든 프록시에서 공통 사용)

WEBAPP_URL = "https://script.google.com/macros/s/AKfycbz3bAOxML5BZSSJcMFM1or5jY8K4NVwliHk_Rbe9jXYVBXbYM05Fl-1bPG1909_38hZ/exec"


# 한글 출력 시 깨지거나 에러 나는 것 방지 (utf-8 인코딩 및 대체 문자 처리)
if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
if hasattr(sys.stderr, "reconfigure"):
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")

# 유틸 함수

# GraphRAG CLI 실행
def _run_graphrag(message, resMethod,paths, resType):
    def decode_output(b: bytes) -> str:
        # subprocess 결과(bytes)를 문자열로 디코딩
        # Windows 환경에서 GraphRAG가 cp949/euc-kr로 출력할 수 있으므로 UTF-8 → CP949 → EUC-KR 순으로 시도
        # 모두 실패하면 UTF-8로 강제 변환 (손실 허용)
        if not b:
            return ""
        for enc in ("utf-8", "cp949", "euc-kr"):
            try:
                return b.decode(enc)
            except UnicodeDecodeError:
                pass
        return b.decode("utf-8", errors="replace")

    # GraphRAG CLI 명령어 구성
    python_command = [
        'graphrag', 'query',
        '--root', paths.GRAPHRAG_ROOT,
        '--response-type', resType,
        '--method', resMethod,
        '--query', message
    ]

    start_time = time.time()

    result = subprocess.run(
        python_command,
        stdout = subprocess.PIPE,
        stderr = subprocess.PIPE,
        env = os.environ.copy(),     # env=os.environ.copy(): 현재 프로세스의 환경변수 상속
        text = False    # text=False: stdout/stderr를 bytes로 받음 (직접 디코딩하기 위해)
    )
    print(f'execution_time : {time.time() - start_time}')

    stdout_text = decode_output(result.stdout)
    stderr_text = decode_output(result.stderr)

    # CLI 오류 (API 키 없음, 인덱스 없음 등)
    if result.returncode != 0:
        raise RuntimeError(stderr_text or stdout_text or 'GraphRAG 실행 오류')

    print(stdout_text)

    # GraphRAG 출력 형식에서 실제 답변 부분만 추출
    match = re.search(r'SUCCESS: (?:Local|Global) Search Response:\s*(.*)', stdout_text, re.DOTALL)
    answer = match.group(1).strip() if match else stdout_text.strip()

    # GraphRAG가 삽입하는 출처 태그 제거
    answer = re.sub(r'\[Data:.*?\]|\[데이터:.*?\]', '', answer)     # 예: "[Data: Sources (1, 2)]", "[데이터: 보고서 (3)]"

    # 마크다운 강조를 평문으로 처리
    answer = re.sub(r'\*+|#+', '', answer)

    answer = answer.strip()
    print(answer)
    return answer.strip()


# 텍스트 → 캘린더 JSON 변환
def _convert_to_calendar_json(text):
    # 자연어 텍스트에서 일정 정보를 추출하여 캘린더 이벤트 JSON으로 변환
    # OpenAI chat completions API를 직접 호출 (GraphRAG 우회, 빠른 응답)
    client = openai.OpenAI(api_key = os.environ.get("GRAPHRAG_API_KEY"))
    try:
        response = client.chat.completions.create(
            model = "gpt-4o-mini",  # gpt-4o-mini 사용: 캘린더 추출은 단순 구조화 작업이므로 저비용 모델로 충분
            response_format = {"type": "json_object"},  # JSON Mode 활성화
            messages = [
                {
                    "role": "system",
                    "content": (
                        "너는 이메일 내용을 분석해서 캘린더 일정을 추출하는 도우미야."
                        "날짜/시간/일정 정보를 추출해서 반드시 JSON으로만 응답해. "
                        "이메일의 제목과 본문을 함께 분석해서 캘린더에 적합한 새로운 일정 제목(title)을 만들어."
                        "메일 제목을 그대로 복사하지 말고, 실제 일정의 목적이 드러나도록 자연스럽고 짧게 작성해."
                        "예를 들면 '회의 안내' 같은 제목이 있더라도, 본문이 캡스톤 발표 회의에 대한 내용이면 title는 '캡스톤 발표 회의'처럼 만들어."
                        "title은 5~20자 정도의 짧고 명확한 한국어로 작성해."
                        "description은 일정과 관련된 핵심 내용을 간단히 넣어"
                        "형식: {\"events\": [{\"title\": \"제목\", \"startTime\": \"2026-02-26 Time 09:00:00\", "
                        "\"endTime\": \"2026-02-26 Time 10:00:00\", \"description\": \"\"}]} "
                        "일정 없으면 {\"events\": []}"
                    )
                },
                {   
                    "role": "user",
                    "content": text
                }
            ]
        )
        return json.loads(response.choices[0].message.content)
    
    except Exception as e:
        # OpenAI API 실패 시 빈 이벤트 반환 (서버 오류 전파 방지)
        print(f"[calendar convert error] {e}")
        return { "events": []}

# 첨부파일 텍스트 요약 (공백/줄바꿈 제외 500자 미만이면 생략)
def _summarize_attachment(text: str, filename: str) -> str:
    pure_len = len(text.replace(" ", "").replace("\n", ""))
    if pure_len < 500:
        return text

    prompt_path = os.path.join("src", "parquet", "prompts", "summarize_attachment.txt")
    with open(prompt_path, "r", encoding="utf-8") as f:
        prompt = f.read().strip()

    client = openai.OpenAI(api_key=os.environ.get("GRAPHRAG_API_KEY"))
    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": prompt},
                {"role": "user", "content": f"파일명: {filename}\n\n{text}"}
            ],
            max_tokens=150
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        print(f"[summarize_attachment error] {e}")
        return text

# PDF 파일에서 텍스트 추출
def _extract_text_from_pdf(file_path):
    text = ""
    try:
        doc = fitz.open(file_path)
        for page in doc:
            text += page.get_text()
        doc.close()
    except Exception as e:
        print(f"[PDF Extract Error] {e}")
    return text

# Word 파일에서 텍스트 추출
def _extract_text_from_docx(file_path):
    text = ""
    try:
        doc = Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        print(f"[Docx Extract Error] {e}")
    return text

# HWP 파일에서 텍스트 추출
def _extract_text_from_hwp(file_path):
    text = ""
    try:
        f = olefile.OleFileIO(file_path)
        dirs = f.listdir()
        sections = [d for d in dirs if "BodyText/Section" in "/".join(d)]
        
        for section in sections:
            stream = f.openstream("/".join(section))
            data = stream.read()
            try:
                # zlib 압축 해제 후 utf-16-le로 디코딩
                decompressed = zlib.decompress(data, -15)
                decoded_text = decompressed.decode("utf-16-le", errors="ignore")
                # 불필요한 제어문자 및 바이너리 찌꺼기 제거
                clean_text = "".join(c for c in decoded_text if c.isalnum() or c in " \n\t.,()[]")
                text += clean_text + "\n"
            except Exception as e:
                print(f"[HWP Decode Error in {section}] {e}")
                
        f.close()
    except Exception as e:
        print(f"[HWP Extract Error] {e}")
    return text

# TXT 파일에서 텍스트 추출 
def _extract_text_from_txt(file_path):
    text = ""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()
    except UnicodeDecodeError:
        try:
            with open(file_path, "r", encoding="cp949") as f:
                text = f.read()
        except Exception as e:
            print(f"[TXT Extract Error] {e}")
    except Exception as e:
        print(f"[TXT Extract Error] {e}")
    return text

# PPTX 파일에서 텍스트 추출
def _extract_text_from_pptx(file_path):
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
    except Exception as e:
        print(f"[PPTX Extract Error] {e}")
    return text

# XLSX 파일에서 텍스트 추출
def _extract_text_from_xlsx(file_path):
    text = ""
    try:
        wb = load_workbook(file_path, data_only=True)
        for ws in wb.worksheets:
            text += f"[Sheet] {ws.title}\n"
            for row in ws.iter_rows(values_only=True):
                row_values = [str(cell) if cell is not None else "" for cell in row]
                # 빈 행은 스킵
                if any(v.strip() for v in row_values):
                    text += " | ".join(row_values) + "\n"
            text += "\n"
    except Exception as e:
        print(f"[XLSX Extract Error] {e}")
    return text

# CSV 파일에서 텍스트 추출
def _extract_text_from_csv(file_path):
    text = ""
    try:
        with open(file_path, "r", encoding="utf-8", newline="") as f:
            reader = csv.reader(f)
            for row in reader:
                row_values = [str(cell) if cell is not None else "" for cell in row]
                text += " | ".join(row_values) + "\n"
    except UnicodeDecodeError:
        try:
            with open(file_path, "r", encoding="cp949", newline="") as f:
                reader = csv.reader(f)
                for row in reader:
                    row_values = [str(cell) if cell is not None else "" for cell in row]
                    text += " | ".join(row_values) + "\n"
        except Exception as e:
            print(f"[CSV Extract Error] {e}")
    except Exception as e:
        print(f"[CSV Extract Error] {e}")
    return text


# 파일명에서 경로/위험 문자 제거
def _sanitize_filename(name: str) -> str:
    name = os.path.basename(name or "attachment.bin").strip()
    name = re.sub(r"[^A-Za-z0-9._-]", "_", name)    # 영숫자, 점, 밑줄, 하이픈만 남기고 나머지는 '_'로 치환
    return name or "attachment.bin"

# attachment payload에서 base64를 받아 서버 로컬에 파일 저장
def _save_attachment_from_base64(file_info: dict, save_dir: str) -> tuple[str, str]:
    original_name = file_info.get("name") or "attachment.bin"
    safe_name = _sanitize_filename(original_name)
    mail_id = str(file_info.get("mail_id") or "no_mail_id")
    data_base64 = file_info.get("data_base64") or ""

    if not data_base64:
        raise ValueError(f"attachment data_base64 missing: {original_name}")

    os.makedirs(save_dir, exist_ok=True)

    ext = os.path.splitext(safe_name)[1].lower()
    unique_name = f"{mail_id}_{uuid.uuid4().hex[:8]}{ext or '.bin'}"
    saved_path = os.path.join(save_dir, unique_name)

    # 혹시 data URL prefix가 붙어오면 제거
    if "," in data_base64 and "base64" in data_base64[:100]:
        data_base64 = data_base64.split(",", 1)[1]

    file_bytes = base64.b64decode(data_base64)

    with open(saved_path, "wb") as f:
        f.write(file_bytes)

    return saved_path, original_name

# 메일 블록에서 'ID: ...' 값을 추출
def _extract_mail_id_from_block(block: str) -> str | None:
    m = re.search(r"^\s*ID:\s*(.+?)\s*$", block, re.MULTILINE)
    return m.group(1).strip() if m else None

# mail_id 기준으로 첨부 텍스트를 각 메일 블록 하단에 삽입한 후 다시 append
def _merge_attachments_into_mail_blocks(content: str, attachment_texts_by_mail: dict[str, list[dict]]) -> str:
    parts = content.split(MAIL_BLOCK_SEP)   # content는 MAIL_BLOCK_SEP 기준으로 메일 블록들이 이어진 문자열이라고 가정
    merged_blocks = []

    for part in parts:
        block = part.strip()
        if not block:
            continue

        # 구분선 복원
        block_text = f"{MAIL_BLOCK_SEP}\n{block}\n{MAIL_BLOCK_SEP}"

        mail_id = _extract_mail_id_from_block(block_text)
        if not mail_id:
            merged_blocks.append(block_text)
            continue

        attachment_entries = attachment_texts_by_mail.get(mail_id, [])
        if not attachment_entries:
            merged_blocks.append(block_text)
            continue

        attachment_section = "\n[첨부 추출 내용]\n"
        for item in attachment_entries:
            attachment_section += f"[File name] {item['name']}\n{item['text']}\n"

        # 블록 하단(마지막 구분선 직전)에 삽입
        insert_pos = block_text.rfind(MAIL_BLOCK_SEP)
        if insert_pos == -1:
            merged_blocks.append(block_text + attachment_section)
        else:
            merged_blocks.append(
                block_text[:insert_pos].rstrip() + "\n\n" +
                attachment_section.rstrip() + "\n" +
                MAIL_BLOCK_SEP
            )

    return "\n".join(merged_blocks) + "\n"

# 텍스트에서 메일별로 구분
def _split_mail_blocks(text):
    parts = text.split(MAIL_BLOCK_SEP) 
    blocks = []

    for p in parts:
        p = p.strip()
        if not p:
            continue
        block = MAIL_BLOCK_SEP + "\n" + p
        if not block.endswith(MAIL_BLOCK_SEP):
            block += "\n" + MAIL_BLOCK_SEP

        blocks.append(block)

    return blocks

# 메일 id들 추출해서 집합으로 반환
def _extract_message_ids(text):
    # re.MULTILINE: ^/$가 각 줄의 시작/끝에 매칭되도록 설정
    return set(re.findall(r"^\s*ID:\s*(.+?)\s*$", text, flags=re.MULTILINE))

# 메일 블록에서 "날짜:" 부분 파싱해서 datetime 객체로 반환
def _extract_block_for_sort(block):
    for line in block.splitlines():
        if line.startswith("날짜:"):
            raw = line.replace("날짜:", "").strip()
            try:
                return datetime.datetime.strptime(raw, "%Y-%m-%d %H:%M:%S")
            except Exception:
                # 날짜 형식이 예상과 다르면 정렬 시 맨 뒤로 정렬
                return datetime.datetime.min
    # 날짜 줄이 없는 경우 정렬 시 맨 뒤로 정렬
    return datetime.datetime.min

# 현재 mail_latest.txt 파일 전체 문자열로 읽어서 반환
def _read_latest_text(paths):
    if not os.path.exists(paths.MAIL_LATEST_PATH):
        return "" # 파일 존재하지 않으면 빈 문자열로 처리
    with open(paths.MAIL_LATEST_PATH, "r", encoding="utf-8") as f:
        return f.read()

# 업데이트 시 생기는 input 폴더 속 새로운 메일 증분 파일 삭제
def _delete_incremental_files(paths):
    os.makedirs(paths.MAIL_DIR, exist_ok=True)

    for name in os.listdir(paths.MAIL_DIR):
        # "inc_"로 시작하고 ".txt"로 끝나는 파일 찾아서 삭제

        if name.startswith("inc_") and name.endswith(".txt"):
            path = os.path.join(paths.MAIL_DIR, name)
            try:
                os.remove(path)
            except Exception as e:
                # 삭제 실패 시 오류 남기고 계속 함
                print(f"[UPLOAD] failed to remove incremental file: {path} / {e}")


# 증분 파일 저장경로 생성
def _build_incremental_path(filename: str, paths) -> str:
    safe_name = _sanitize_filename(filename or "") # 경로 탐색 공격 등 방지용 정제
    # 정제 후에도 "inc_"로 시작하지 않으면 시간 기반 파일명으로 대체
    if not safe_name.startswith("inc_"):
        safe_name = f"inc_{datetime.datetime.now().strftime('%Y-%m-%d_%H%M%S')}.txt"
    return os.path.join(paths.MAIL_DIR, safe_name)

# json 파일 읽어서 dict로 파싱 후 반환
def _read_json_file(path):
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)


# 이름+메일주소 형식에서 이름과 메일주소 분리하여 반환
def _parse_contact(raw: str) -> tuple[str, str]:
        m = re.search(r"^(.*?)\s*<([^>]+)>", raw.strip())
        if m:
            name  = m.group(1).strip().strip('"')
            email = m.group(2).strip().lower()
        else:
            name  = ""
            email = raw.strip().lower()
        return name, email

# 메일 블록에서 특정 필드 값 추출
def _extract_field(block: str, label: str) -> str:
        m = re.search(rf"^{label}:\s*(.+)$", block, re.MULTILINE)
        return m.group(1).strip() if m else ""

# 메일 발신 수신 횟수 계정별로 저장
def _save_mail_contact_stats(blocks: list[str],paths, mode: str = "rewrite"):
    
    # 새로운 메일만 추가된 거라 이미 횟수 저장한 json 파일이 존재할 때 
    if mode == "append" and os.path.exists(paths.MAIL_STATICS_PATH):
        with open(paths.MAIL_STATICS_PATH, "r", encoding="utf-8") as f:
            stats = json.load(f)
    else: # 전체 갱신 모드일 때 빈 딕셔너리로 초기화해서 새로 횟수 셈
        stats = {}
    # 송수신 횟수 누적
    def add(name: str, email: str, direction: str):
        if not email or email in ("-", ""):
            return
        # 이메일 처음 등장하면 name, sent, received 초기화
        stats.setdefault(email, {"name": name, "sent": 0, "received": 0})
        # 이름이 있을 때 덮어씀
        if name:
            stats[email]["name"] = name
        stats[email][direction] += 1
    # 블록 순회하며 횟수 집계
    for block in blocks:
        direction = _extract_field(block, "구분") # 발신 또는 수신
        from_raw  = _extract_field(block, "발신인") # 발신인 원문
        to_raw    = _extract_field(block, "수신인") # 수신인 원문 

        if direction == "발신":
            # 수신인 여러명이면 ,로 구분
            for addr in to_raw.split(","):
                name, email = _parse_contact(addr)
                add(name, email, "sent")
        elif direction == "수신":
            name, email = _parse_contact(from_raw)
            add(name, email, "received")

    # json 파일에 저장
    os.makedirs(os.path.dirname(paths.MAIL_STATICS_PATH), exist_ok=True)    
    with open(paths.MAIL_STATICS_PATH, "w", encoding="utf-8") as f:
        json.dump(stats, f, ensure_ascii=False, indent=2) # indent=2 : 사람이 읽기 쉽게 들여쓰기 적용

    print(f"[STATS] ({mode}) 계정 {len(stats)}개 집계 완료 → {paths.MAIL_STATICS_PATH}")
    
# 인덱스 여부 확인
def _is_index_ready(paths):

    #graph_path = os.path.join(paths.GRAPHRAG_ROOT, "output", "graph.graphml")
    stats_path = os.path.join(paths.GRAPHRAG_ROOT, "output", "stats.json")

    try:
         # 동기화된 메일 텍스트, graphml 파일, 인덱싱 통계 파일이 존재하는지 확인
        required_paths = [paths.MAIL_LATEST_PATH, stats_path]

        for path in required_paths:
            if not os.path.exists(path):
                print(f"[INDEX READY] missing: {path}")
                return False
            if os.path.getsize(path) == 0:
                print(f"[INDEX READY] empty file: {path}")
                return False

        # stats.json JSON 파싱 시도해서 파일 손상 여부 확인
        _read_json_file(stats_path)
        return True

    except Exception as e:
        # 예상치 못한 오류는 인덱스 불완전으로 판단함
        print(f"[INDEX READY] invalid index state: {e}")
        return False

# 엔드포인트: POST /extract-calendar
@app.route('/extract-calendar', methods=['POST'])
def extract_calendar():     # 이메일 제목 + 본문에서 일정 이벤트를 추출하여 반환
    data = request.json or {}
    subject = data.get('subject', '')
    body = data.get('body', '')
    result = _convert_to_calendar_json(f"제목: {subject}\n\n{body}")    # 제목과 본문을 합쳐 컨텍스트 제공
    return jsonify(result)

# 엔드포인트: POST /run-query-async
@app.route('/run-query-async', methods=['POST'])    # GraphRAG 쿼리를 백그라운드 스레드에서 비동기 실행하고 Job ID를 즉시 반환
def run_query_async():
    data = request.json or {}

    print("[DEBUG] content_type =", request.content_type)
    print("[DEBUG] raw body =", request.data)
    print("[DEBUG] parsed data =", data)

    if data is None:
        return jsonify({'error': 'JSON 본문을 읽지 못했습니다.'}), 400

    message = request.json.get('message', '')
    resMethod = request.json.get('resMethod', 'local')
    resType = request.json.get('resType', 'text')
    gmail_id = data.get('gmail_id', '').strip()

    if not str(message).strip():
        return jsonify({'error': 'message가 비어있습니다.'}), 400

    if not gmail_id:
        return jsonify({'error': 'gmail_id가 비어있습니다.'}), 400

    print("[DEBUG] message =", repr(message))
    print("[DEBUG] gmail_id =", repr(gmail_id))
    
    # uuid4: 랜덤 UUID 생성. [:8]로 앞 8자리만 사용 (충돌 가능성 낮고 가독성 좋음)
    job_id = str(uuid.uuid4())[:8]
    create_job(job_id, job_type="query")
    update_job(job_id, status="pending", result=None, resType=resType)

    def _worker():  # 백그라운드 스레드에서 실행되는 실제 작업 함수
        try:

            paths = UserPaths(BASE_DIR, gmail_id)
            env = os.environ.copy()
            env["GMAIL_ID"] = gmail_id

            # 한국어 응답 강제 (GraphRAG 기본 응답이 영어일 경우 대비)
            full_message = message + " 영어 말고 한국어로 답변해줘."
            answer = _run_graphrag(full_message, resMethod, paths, resType )

            if resType.lower() == "calendar":
                # 캘린더 타입: GraphRAG 텍스트 답변을 다시 OpenAI로 구조화
                result = json.dumps(_convert_to_calendar_json(answer), ensure_ascii=False)
            else:
                result = answer

            update_job(job_id, status="done", result=result)

        except Exception as e:
            update_job(job_id, status="error", result=str(e))

    # daemon=True: 메인 프로세스 종료 시 스레드도 함께 종료
    threading.Thread(target=_worker, daemon=True).start()
    return jsonify({"jobId": job_id})

# 엔드포인트: GET /job-status/<job_id>
@app.route('/job-status/<job_id>', methods=['GET'])
def job_status(job_id):     # 비동기 Job의 현재 상태와 결과를 반환

    job = get_job(job_id)
    if not job:
        return jsonify({"status": "not_found"}), 404

    if job["status"] == "done" and job["resType"].lower() == "calendar":
        try:
            # 캘린더 결과는 JSON 문자열로 저장되어 있으므로 파싱 후 반환
            return jsonify({"status": "done", "data": json.loads(job["result"])})
        except Exception:
            return jsonify({"status": "done", "data": {"events": []}})

    # text 타입: result 필드에 문자열 그대로 반환
    return jsonify({"status": job["status"], "result": job["result"] or ""})

# 엔드포인트: POST /run-query  (동기 버전, 디버깅/단순 클라이언트용)
@app.route('/run-query', methods=['POST'])
def run_query():    # GraphRAG 쿼리를 동기 방식으로 실행하고 결과를 즉시 반환
    message = request.json.get('message', '')
    resMethod = request.json.get('resMethod', 'local')
    resType = request.json.get('resType', 'text')

    print(f'message: {message}')
    print(f'resMethod: {resMethod}')
    print(f'resType: {resType}')

    if not str(message).strip():
        return jsonify({'error': 'message가 비어있습니다.'}), 400

    message += " 영어 말고 한국어로 답변해줘."

    try:
        answer = _run_graphrag(message, resMethod, resType)
    except RuntimeError as e:
        return jsonify({'error': str(e)}), 500

    if resType.lower() == "calendar":
        return jsonify(_convert_to_calendar_json(answer))

    return jsonify({'result': answer})

# 엔드포인트: POST /upload
@app.route("/upload", methods=["POST"])
def upload():
    # 1) 데이터 수신
    data = request.json or {}
    filename = data.get("filename") or f"mail_{int(time.time())}.txt"
    content = data.get("content") or ""
    attachments = data.get("attachment") or []
    requested_mode = data.get("syncmode", "append")
    gmail_id = (data.get("gmail_id") or "").strip().lower()

    paths = UserPaths(BASE_DIR, gmail_id) # 각 유저별 고유경로 설정

    if not str(content).strip():
        return jsonify({"ok": False, "error": "content가 비어있습니다."}), 400
    if not gmail_id:
        return jsonify({"ok": False, "error": "gmail_id가 비어있습니다."}), 400
    
    print("user gmail id =", gmail_id)
    # append인데 기존 인덱스가 없으면 rewrite로 전환
    fallback_to_rewrite = False
    sync_mode = requested_mode

    # 새로운 메일 추가 모드이지만 인덱싱이 되어있지 않으면 인덱싱 모드로 전환
    if requested_mode == "append" and not _is_index_ready(paths):
        print("[UPLOAD] index not ready -> fallback to rewrite")
        sync_mode = "rewrite"
        fallback_to_rewrite = True

    # 2) 저장 디렉토리 준비


    os.makedirs(paths.MAIL_DIR, exist_ok=True)

    # rewrite면 기존 첨부파일 먼저 전부 삭제
    if sync_mode == "rewrite":
        if os.path.exists(paths.ATTACHMENT_DIR):
            shutil.rmtree(paths.ATTACHMENT_DIR)
            print(f"[CLEAN] attachment 폴더 초기화 완료: {paths.ATTACHMENT_DIR}")

    file_path = os.path.join(paths.MAIL_DIR, filename)

    # 3) 원본 메일 텍스트 저장
    with open(file_path, "w", encoding="utf-8") as f:
        f.write(content)

    extracted_count = 0
    failed_attachments = []
    saved_attachment_paths = []
    attachment_texts_by_mail: dict[str, list[dict]] = {}

    # 5) 첨부 저장 + 텍스트 추출 + mail_id별 묶기
    if attachments:
        extracted_full_text = f"\n\n{MAIL_BLOCK_SEP}\n"
        extracted_full_text += "[System] attachment data extract section\n"

        for file_info in attachments:
            f_name = file_info.get("name") or "attachment.bin"  
            mime = (file_info.get("mime") or "").lower()        
            mail_id = str(file_info.get("mail_id") or "").strip()

            try:
                # base64 → 서버 로컬 파일 저장
                saved_path, original_name = _save_attachment_from_base64(file_info, paths.ATTACHMENT_DIR)
                saved_attachment_paths.append(saved_path)

                ext = os.path.splitext(original_name)[-1].lower()
                file_text = ""

                # MIME type 제한
                if ext == ".pdf" or mime in ("application/pdf", "application/haansoftpdf"):     
                    file_text = _extract_text_from_pdf(saved_path)
                elif ext == ".docx" or mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                    file_text = _extract_text_from_docx(saved_path)
                elif ext == ".hwp" or mime in ("application/x-hwp", "application/haansofthwp"): # 추가
                    file_text = _extract_text_from_hwp(saved_path)
                elif ext == ".txt" or mime == "text/plain":
                    file_text = _extract_text_from_txt(saved_path)
                elif ext == ".pptx" or mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                    file_text = _extract_text_from_pptx(saved_path)
                elif ext == ".xlsx" or mime == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                    file_text = _extract_text_from_xlsx(saved_path)
                elif ext == ".csv" or mime in ("text/csv", "application/csv"):
                    file_text = _extract_text_from_csv(saved_path)
                else:
                    failed_attachments.append({
                        "name": original_name,
                        "reason": f"unsupported type: ext={ext}, mime={mime}"
                    })
                    continue

                if file_text and file_text.strip():
                    if mail_id:
                        # 텍스트만 저장, 요약은 백그라운드에서 처리
                        attachment_texts_by_mail.setdefault(mail_id, []).append({
                            "name": original_name,
                            "text": file_text.strip()
                        })
                    else:
                        failed_attachments.append({
                            "name": original_name,
                            "reason": "mail_id missing"
                        })
                    extracted_count += 1
                else:
                    failed_attachments.append({
                        "name": original_name,
                        "reason": "text extraction returned empty"
                    })
                    continue

            except Exception as e:
                failed_attachments.append({
                    "name": f_name,
                    "reason": str(e)
                })
                print(f"[UPLOAD][ATTACHMENT ERROR] {f_name}: {e}")  

    # 7) 파이프라인 실행
    print(f"[UPLOAD] Received filename: {filename}")
    print(f"[UPLOAD] Content length: {len(content)}")
    print(f"[UPLOAD] Attachment count received: {len(attachments)}")
    print(f"[UPLOAD] Attachment extracted count: {extracted_count}")
    print(f"[UPLOAD] Requested mode: {requested_mode}")
    print(f"[UPLOAD] Actual mode: {sync_mode}")
    print("[UPLOAD] cwd:", os.getcwd())

    added_count = 0 # 저장된 메일 블록 수
    skipped_count = 0 # 건너뛰는 메일 수
    saved_mail_path = "" # 최종 저장 파일 경로

    if sync_mode == "rewrite":

        # 첨부 병합 없이 메일 본문만 저장 (첨부 요약+병합은 백그라운드에서 처리)
        _delete_incremental_files(paths)
        with open(paths.MAIL_LATEST_PATH, "w", encoding="utf-8") as f:
            f.write(content.rstrip() + "\n")
        saved_mail_path = paths.MAIL_LATEST_PATH

        added_count = len(_split_mail_blocks(content))

        added_count = len(_split_mail_blocks(content))
        _save_mail_contact_stats(_split_mail_blocks(content), paths, mode="rewrite")

    # 새 메일만 추가 append 모드
    else:

        # 기존 mail_latest.txt에서 인덱싱된 메일 ID 추출해서 중복 방지
        existing_text = _read_latest_text(paths)
        existing_ids = _extract_message_ids(existing_text)
        new_blocks = _split_mail_blocks(content)
        append_blocks = []

        for block in new_blocks:
            msg_id = _extract_mail_id_from_block(block)

            if not msg_id:
                skipped_count += 1
                continue

            if msg_id in existing_ids:
                skipped_count += 1
                continue

            if msg_id in attachment_texts_by_mail:
                block = _merge_attachments_into_mail_blocks(
                    block,
                    {msg_id: attachment_texts_by_mail[msg_id]}
                ).strip()

            append_blocks.append(block.strip())
            existing_ids.add(msg_id)

        added_count = len(append_blocks)

        if append_blocks:
            append_blocks.sort(key=_extract_block_for_sort, reverse=True)
            inc_content = "\n\n".join(append_blocks).strip() + "\n"

            # 시간 기반 파일명으로 증분파일 저장
            inc_path = _build_incremental_path(filename, paths)


            with open(inc_path, "w", encoding="utf-8") as f:
                f.write(inc_content)

            updated_content = inc_content + "\n" + existing_text
            with open(paths.MAIL_LATEST_PATH, "w", encoding="utf-8") as f:
                f.write(updated_content.strip() + "\n")

            saved_mail_path = inc_path

            _save_mail_contact_stats(append_blocks, paths, mode="append")

        else:
            saved_mail_path = ""

    print("[UPLOAD] added:", added_count)
    print("[UPLOAD] skipped:", skipped_count)
    if saved_mail_path:
        print("[UPLOAD] saved mail path:", os.path.abspath(saved_mail_path))

    # GraphRAG 파이프라인을 백그라운드에서 실행
    job_id = str(uuid.uuid4())[:8] # 작업 구분용, 앞 8자리만 잘라서 사용

    if sync_mode == "rewrite":
        create_job(job_id, job_type="index") # 새로운 작업을 생성 (타입: index = 전체 재생성)
        update_job(job_id, message="업로드 완료, 그래프 파이프라인 시작") # 작업 상태 메시지 업데이트 (로그에서 확인용)

    else:
        create_job(job_id, job_type="update") # 증분 업데이트 작업 등록
        update_job(job_id, message="업로드 완료, 그래프 업데이트 파이프라인 시작") # 작업 상태 메시지 업데이트

    env = os.environ.copy() # os.environ = 프로세스의 환경변수들을 담고 있는 객체, 모든 프로세스의 환경을 통일하기 위함
    env["PYTHONUNBUFFERED"] = "1" # 실시간 로그를 출력하기 위함


    if sync_mode == "rewrite": # 전체 갱신할 때

        update_dir = os.path.join(paths.GRAPHRAG_ROOT, "update_output") # 이전에 증분 결과 있으면 폴더 삭제
        if os.path.exists(update_dir): 
            shutil.rmtree(update_dir)
            print(f"[CLEAN] update_output 삭제 완료: {update_dir}")
        else:
            print(f"[CLEAN] update_output 없음: {update_dir}")

        start_graph_pipeline_background(job_id,paths, env, attachment_texts_by_mail) # GraphRAG 파이프라인 함수 실행

    else:
        start_graph_update_pipeline_background(job_id,paths, env)

    return jsonify({
            "ok": True,
            "requested_mode": requested_mode,
            "job_id":job_id,
            "actual_mode": sync_mode,
            "fallback_to_rewrite": fallback_to_rewrite,
            "latest_path": os.path.abspath(paths.MAIL_LATEST_PATH),
            "saved_mail_path": os.path.abspath(saved_mail_path) if saved_mail_path else "",
            "attachment_dir": os.path.abspath(paths.ATTACHMENT_DIR),
            "content_length": len(content),
            "added_count": added_count,
            "skipped_count": skipped_count,
            "attachment_received_count": len(attachments),
            "attachment_extracted_count": extracted_count,
            "failed_attachments": failed_attachments,
        })

# 엔드포인트: GET /graph-data
@app.route("/graph-data", methods=["GET", "OPTIONS"])  # OPTIONS 추가 (CORS 프리플라이트)
def graph_data():
    if request.method == "OPTIONS":
        return "", 200

    gmail_id = (request.args.get("gmail_id") or "").strip().lower()

    if not gmail_id:
        return jsonify({"ok": False, "error": "gmail_id가 비어있습니다."}), 400

    paths = UserPaths(BASE_DIR, gmail_id)  # 각 유저별 고유경로 설정

    if not os.path.exists(paths.GRAPH_JSON_PATH):
        return jsonify({"nodes": [], "edges": [], "error": "graph json not found"}), 200

    try:
        with open(paths.GRAPH_JSON_PATH, "r", encoding="utf-8") as f:
            data = json.load(f)
        print(f"[GRAPH-DATA] 반환: {len(data.get('nodes', []))} 노드")  # 로그 추가
        return jsonify(data)
    except Exception as e:
        print(f"[GRAPH-DATA] 에러: {e}")
        return jsonify({"nodes": [], "edges": [], "error": str(e)}), 500

# 엔드포인트: GET /graph-view
@app.route("/graph-view", methods=["GET"])
# 브라우저에서/graph-view 접속하면 graph_view.html 반환
def graph_view():
    return send_from_directory(
        os.path.join(os.path.dirname(__file__), "json"), # src/json/ 폴더
        "graph_view.html"
    )

# 공유 그래프 렌더링 함수 (graph_view.html 과 dashboard 양쪽에서 사용)
@app.route('/graph-render.js')
def graph_render_js():
    return send_from_directory(
        os.path.join(os.path.dirname(__file__), "json"),
        "graph-render.js"
    )

# 엔드포인트: GET /index-status
@app.route("/index-status", methods=["GET"])
def index_status():     # GraphRAG 인덱싱 완료 여부 반환
    return jsonify({ "indexed": _is_index_ready() })

# 엔드포인트: GET /dashboard/ (Gentella 웹앱 서빙)
@app.route('/dashboard/', defaults={'path': 'production/index.html'})
@app.route('/dashboard/<path:path>')
def dashboard(path):
    dist_dir = os.path.join(os.path.dirname(__file__), 'apps-script', 'web', 'dist')
    # /dashboard/index2.html 요청 → production/index2.html로 매핑
    if not path.startswith('production/') and path.endswith('.html'):
        path = 'production/' + path
    return send_from_directory(dist_dir, path)

# dist 루트 정적 파일 서빙 (assets, js, fonts)
@app.route('/assets/<path:path>')
def static_assets(path):
    dist_dir = os.path.join(os.path.dirname(__file__), 'apps-script', 'web', 'dist', 'assets')
    return send_from_directory(dist_dir, path)

@app.route('/js/<path:path>')
def static_js(path):
    dist_dir = os.path.join(os.path.dirname(__file__), 'apps-script', 'web', 'dist', 'js')
    return send_from_directory(dist_dir, path)

@app.route('/fonts/<path:path>')
def static_fonts(path):
    dist_dir = os.path.join(os.path.dirname(__file__), 'apps-script', 'web', 'dist', 'fonts')
    return send_from_directory(dist_dir, path)

# 엔드포인트: POST /calendar-events (Apps Script 캘린더 프록시)
@app.route('/calendar-events', methods=['POST'])
def calendar_events():
    data = request.json or {}
    res = requests.post(WEBAPP_URL, json=data, allow_redirects=True)
    print("[calendar] status:", res.status_code)
    print("[calendar] response:", res.text[:500])
    try:
        return jsonify(res.json())
    except Exception:
        return jsonify({"events": [], "error": res.text[:200]}), 200

# 엔드포인트: POST /labels-proxy (Apps Script 라벨 프록시)
@app.route('/labels-proxy', methods=['POST'])
def labels_proxy():
    data = request.json or {}
    try:
        res = requests.post(WEBAPP_URL, json=data, allow_redirects=True)
        print("[labels] status:", res.status_code)
        try:
            return jsonify(res.json())
        except Exception:
            return jsonify({"ok": False, "error": res.text[:200]}), 200
    except Exception as e:
        return jsonify({"ok": False, "error": str(e)}), 500

import urllib.request

@app.route('/dashboard/marker-icon.png')
def marker_icon():
    url = 'https://cdnjs.cloudflare.com/ajax/libs/leaflet/1.9.4/images/marker-icon.png'
    with urllib.request.urlopen(url) as r:
        data = r.read()
    from flask import Response
    return Response(data, mimetype='image/png')

@app.route('/dashboard/marker-shadow.png')
def marker_shadow():
    url = 'https://cdnjs.cloudflare.com/ajax/libs/leaflet/1.9.4/images/marker-shadow.png'
    with urllib.request.urlopen(url) as r:
        data = r.read()
    from flask import Response
    return Response(data, mimetype='image/png')

# 서버 진입점
if __name__ == '__main__':
    # host='0.0.0.0': 모든 네트워크 인터페이스에서 수신 (localhost 외부 접근 허용)
    # port=80: 표준 HTTP 포트. Linux에서는 root 권한 필요 (또는 포트포워딩 사용)
    # debug=False: 운영 환경 설정. True로 바꾸면 코드 변경 시 자동 재시작, 에러 상세 표시
    app.run(host='0.0.0.0', port=80, debug=False)